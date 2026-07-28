"""Office documents, workbooks and decks.

The property under test everywhere here: **a spreadsheet must stay a
spreadsheet**, and a figure EMEFA states must be one she computed, never one
she guessed. Files are opened back with the same libraries so the assertions
are about what a user would actually see.
"""

import io

import pytest
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation

from emefa.config import Settings
from emefa.domain.agent import AgentStep
from emefa.domain.documents import ARTIFACT_MIME_TYPES, DocumentStore
from emefa.domain.office import (
    Block,
    BlockKind,
    Brand,
    Chart,
    ChartKind,
    Column,
    ColumnFormat,
    DeckSpec,
    DocumentKind,
    DocumentSpec,
    OfficeError,
    Sheet,
    Slide,
    SlideLayout,
    WorkbookSpec,
)
from emefa.infrastructure.office_native import (
    NativeOfficeProvider,
    _evaluate_row_formula,
)
from emefa.main import create_app

PROVIDER = NativeOfficeProvider()
BRAND = Brand(company_name="Horizon SARL", contact="contact@horizon.tg")


# ── Word ──────────────────────────────────────────────────────────────────


def quote_spec(**overrides):
    defaults = dict(
        kind=DocumentKind.QUOTE,
        title="Devis 2026-042",
        reference="DV-2026-042",
        recipient="Clinique du Lac",
        brand=BRAND,
        blocks=(
            Block(BlockKind.HEADING, "Prestations", level=1),
            Block(
                BlockKind.TABLE,
                columns=("Désignation", "Qté", "Montant"),
                rows=(("Installation", "2", "2 400 000"),),
                numeric_last_column=True,
            ),
            Block(BlockKind.SIGNATURE, items=("Le client", "Horizon SARL")),
        ),
    )
    return DocumentSpec(**{**defaults, **overrides})


def test_a_quote_carries_what_makes_it_a_quote():
    built = PROVIDER.build_document(quote_spec())
    document = DocxDocument(io.BytesIO(built.data))
    text = "\n".join(paragraph.text for paragraph in document.paragraphs)

    assert built.extension == "docx"
    assert "Devis 2026-042" in text
    assert "DV-2026-042" in text
    assert "Clinique du Lac" in text
    assert "Horizon SARL" in text
    # A table with its header, and a signature block.
    assert document.tables[0].rows[0].cells[0].text == "Désignation"
    assert "Signature" in document.tables[-1].cell(0, 0).text


def test_headers_footers_and_page_numbers_are_fields_not_text():
    """A page number written as text is wrong on every page but one."""
    built = PROVIDER.build_document(quote_spec())
    document = DocxDocument(io.BytesIO(built.data))
    section = document.sections[0]

    assert "Horizon SARL" in section.header.paragraphs[0].text
    footer_xml = section.footer.paragraphs[0]._p.xml
    assert "PAGE" in footer_xml
    assert "fldChar" in footer_xml


def test_a_table_of_contents_is_a_field_and_says_so():
    built = PROVIDER.build_document(quote_spec(table_of_contents=True))
    document = DocxDocument(io.BytesIO(built.data))
    body = document.element.body.xml

    assert "TOC" in body and "fldChar" in body
    # The user must know it fills in on open, or they will think it is broken.
    assert any("sommaire se remplit" in note.lower() for note in built.warnings)


def test_every_block_kind_renders():
    built = PROVIDER.build_document(
        quote_spec(
            kind=DocumentKind.REPORT,
            blocks=(
                Block(BlockKind.HEADING, "Contexte", level=1),
                Block(BlockKind.PARAGRAPH, "Le marché progresse."),
                Block(BlockKind.BULLETS, items=("Point un", "Point deux")),
                Block(BlockKind.NUMBERED, items=("Étape un",)),
                Block(BlockKind.FIELDS, fields=(("Client", "Horizon"), ("Montant", "4 M"))),
                Block(BlockKind.PAGE_BREAK),
                Block(BlockKind.TABLE, columns=("A",), rows=(("1",),)),
            ),
        )
    )
    document = DocxDocument(io.BytesIO(built.data))
    text = "\n".join(paragraph.text for paragraph in document.paragraphs)
    assert "Le marché progresse." in text
    assert "Point deux" in text
    assert "Étape un" in text
    assert len(document.tables) == 2  # the fields block and the table block


def test_an_untitled_document_is_refused():
    with pytest.raises(OfficeError):
        PROVIDER.build_document(quote_spec(title="   "))


# ── Excel ─────────────────────────────────────────────────────────────────


def devis_workbook():
    return WorkbookSpec(
        title="Devis chiffré",
        brand=BRAND,
        sheets=(
            Sheet(
                name="Devis",
                caption="Clinique du Lac — juillet 2026",
                columns=(
                    Column("Désignation"),
                    Column("Qté", ColumnFormat.NUMBER),
                    Column("PU", ColumnFormat.MONEY),
                    Column("Montant", ColumnFormat.MONEY, formula="=B{row}*C{row}", total=True),
                ),
                rows=(
                    ("Installation", 2, 1_200_000, None),
                    ("Maintenance", 12, 80_000, None),
                ),
                chart=Chart(ChartKind.BAR, "Répartition", "Désignation", ("Qté",)),
            ),
        ),
    )


def test_a_workbook_stays_a_workbook():
    """Values written where formulas belong die the moment someone edits a
    quantity."""
    built = PROVIDER.build_workbook(devis_workbook())
    workbook = load_workbook(io.BytesIO(built.data))
    sheet = workbook["Devis"]

    # Caption on row 1, header on row 3, data from row 4.
    assert sheet["A1"].value.startswith("Clinique du Lac")
    assert sheet["A3"].value == "Désignation"
    assert sheet["D4"].value == "=B4*C4"
    assert sheet["D5"].value == "=B5*C5"
    assert sheet["D6"].value == "=SUM(D4:D5)"
    assert sheet["C4"].number_format == "#,##0"
    assert sheet.freeze_panes == "A4"
    assert sheet.auto_filter.ref == "A3:D5"
    assert len(sheet._charts) == 1


def test_the_total_of_a_formula_column_is_computed_not_guessed():
    """Summing the raw cells of a formula column gives zero, and an assistant
    announcing "total : 0" on a quote is worse than one that says it cannot
    tell."""
    built = PROVIDER.build_workbook(devis_workbook())
    assert built.computed["Devis.Montant"] == pytest.approx(2 * 1_200_000 + 12 * 80_000)
    assert built.warnings == ()


def test_an_unevaluable_formula_is_admitted_rather_than_reported_as_zero():
    built = PROVIDER.build_workbook(
        WorkbookSpec(
            title="Complexe",
            sheets=(
                Sheet(
                    name="F",
                    columns=(
                        Column("A", ColumnFormat.NUMBER),
                        Column("B", ColumnFormat.MONEY, formula="=SUM(A1:A9)", total=True),
                    ),
                    rows=((1, None),),
                ),
            ),
        )
    )
    assert "F.B" not in built.computed
    assert any("apparaîtra à l'ouverture" in note for note in built.warnings)


@pytest.mark.parametrize(
    ("formula", "row", "expected"),
    [
        ("=B{row}*C{row}", ("x", 2, 100.0), 200.0),
        ("=B{row}*C{row}*1.18", ("x", 2, 100.0), 236.0),
        ("=B{row}-C{row}", ("x", 10, 3), 7.0),
        ("=B{row}+C{row}*2", ("x", 10, 3), 16.0),  # precedence, not left to right
        ("=B{row}/C{row}", ("x", 10, 4), 2.5),
        ("=B{row}/C{row}", ("x", 10, 0), None),  # division by zero is not a number
        ("=SUM(A1:A9)", ("x", 1), None),  # left to Excel
        ("=B{row}*Z{row}", ("x", 2), None),  # reference beyond the row
        ("=IF(B{row}>0,1,0)", ("x", 2), None),
    ],
)
def test_row_formula_evaluation(formula, row, expected):
    result = _evaluate_row_formula(formula, row)
    if expected is None:
        assert result is None
    else:
        assert result == pytest.approx(expected)


def test_sheet_names_excel_would_reject_are_repaired():
    built = PROVIDER.build_workbook(
        WorkbookSpec(
            title="X",
            sheets=(Sheet(name="Ventes/2026[T1]", columns=(Column("A"),), rows=(("1",),)),),
        )
    )
    assert load_workbook(io.BytesIO(built.data)).sheetnames == ["Ventes-2026-T1-"]


def test_a_chart_referencing_a_missing_column_is_reported_not_silently_dropped():
    built = PROVIDER.build_workbook(
        WorkbookSpec(
            title="X",
            sheets=(
                Sheet(
                    name="F",
                    columns=(Column("A"),),
                    rows=(("1",),),
                    chart=Chart(ChartKind.BAR, "T", "Inconnue", ("A",)),
                ),
            ),
        )
    )
    assert any("Graphique ignoré" in note for note in built.warnings)


def test_an_empty_workbook_is_refused():
    with pytest.raises(OfficeError):
        PROVIDER.build_workbook(WorkbookSpec(title="Vide"))


# ── PowerPoint ────────────────────────────────────────────────────────────


def test_a_deck_renders_every_layout():
    built = PROVIDER.build_deck(
        DeckSpec(
            title="Comité de direction",
            subtitle="Juillet 2026",
            brand=BRAND,
            slides=(
                Slide(SlideLayout.TITLE),
                Slide(SlideLayout.SECTION, "Commercial"),
                Slide(SlideLayout.BULLETS, "Faits marquants", bullets=("Un", "Deux")),
                Slide(
                    SlideLayout.TABLE,
                    "Pipeline",
                    columns=("Client", "Étape"),
                    rows=(("Horizon", "Devis"),),
                ),
                Slide(
                    SlideLayout.CHART,
                    "Ventes",
                    points=(("Jan", 3.0), ("Fév", 5.0)),
                    chart=Chart(ChartKind.BAR, "Ventes"),
                ),
                Slide(SlideLayout.QUOTE, "Le mot du DG", subtitle="On accélère."),
                Slide(SlideLayout.CLOSING, "Merci", notes="Rappeler la clinique."),
            ),
        )
    )
    deck = Presentation(io.BytesIO(built.data))

    assert built.extension == "pptx"
    assert len(deck.slides) == 7
    titles = [slide.shapes.title.text for slide in deck.slides if slide.shapes.title]
    assert "Faits marquants" in titles
    assert any(shape.has_table for slide in deck.slides for shape in slide.shapes)
    assert any(shape.has_chart for slide in deck.slides for shape in slide.shapes)
    assert deck.slides[-1].notes_slide.notes_text_frame.text == "Rappeler la clinique."


def test_a_chart_slide_without_data_is_reported():
    built = PROVIDER.build_deck(
        DeckSpec(title="X", slides=(Slide(SlideLayout.CHART, "Vide"),))
    )
    assert any("sans données" in note for note in built.warnings)


def test_an_empty_deck_is_refused():
    with pytest.raises(OfficeError):
        PROVIDER.build_deck(DeckSpec(title="Vide"))


# ── storage ───────────────────────────────────────────────────────────────


def test_artifacts_of_every_format_live_side_by_side(tmp_path):
    store = DocumentStore(tmp_path / "emefa.db")
    word = store.create("Compte rendu", "Contenu")
    book = store.save_artifact(b"fake-xlsx", "xlsx", "Budget 2026", "workbook")
    deck = store.save_artifact(b"fake-pptx", "pptx", "Comité", "deck")

    assert store.describe(book["document_id"])["content_type"] == ARTIFACT_MIME_TYPES["xlsx"]
    assert store.describe(book["document_id"])["filename"].endswith(".xlsx")
    assert store.describe(deck["document_id"])["kind"] == "deck"
    # A Word document created before spreadsheets existed still reads.
    assert store.describe(word["document_id"])["title"] == "Compte rendu"

    listed = {item["document_id"] for item in store.list()}
    assert listed == {word["document_id"], book["document_id"], deck["document_id"]}

    with pytest.raises(ValueError):
        store.save_artifact(b"x", "exe", "Malveillant")


# ── through the agent ─────────────────────────────────────────────────────


def build_app(tmp_path):
    class Brain:
        async def think(self, history, tools):
            return AgentStep(answer="Ok.")

    return create_app(
        Settings(
            enrollment_code="CODE-SECRET",
            database_path=tmp_path / "office.db",
            cookie_secure=False,
        ),
        brain=Brain(),
    )


def test_the_agent_produces_real_office_files(tmp_path):
    app = build_app(tmp_path)
    tools = app.state.agent.tools
    app.state.profiles.update_business({"company_name": "Horizon SARL"})

    document = tools.get("office_document").handler(
        {
            "kind": "quote",
            "title": "Devis 2026-042",
            "recipient": "Clinique du Lac",
            "blocks": [
                {"kind": "heading", "text": "Prestations"},
                {
                    "kind": "table",
                    "columns": ["Désignation", "Montant"],
                    "rows": [["Installation", "2 400 000"]],
                },
            ],
        }
    )
    assert document["document"]["filename"].endswith(".docx")
    # The company's own name, taken from the profile rather than invented.
    stored = DocxDocument(str(app.state.documents.get(document["document"]["document_id"])))
    assert "Horizon SARL" in "\n".join(p.text for p in stored.paragraphs)

    workbook = tools.get("office_spreadsheet").handler(
        {
            "title": "Suivi des ventes",
            "sheets": [
                {
                    "name": "Ventes",
                    "columns": [
                        {"header": "Mois"},
                        {"header": "Qté", "format": "number"},
                        {"header": "PU", "format": "money"},
                        {
                            "header": "Montant",
                            "format": "money",
                            "formula": "=B{row}*C{row}",
                            "total": True,
                        },
                    ],
                    "rows": [["Janvier", 3, 150000, None], ["Février", 5, 150000, None]],
                }
            ],
        }
    )
    assert workbook["computed"]["Ventes.Montant"] == 1_200_000
    assert workbook["document"]["filename"].endswith(".xlsx")

    deck = tools.get("office_deck").handler(
        {
            "title": "Comité",
            "slides": [
                {"layout": "title"},
                {"layout": "bullets", "title": "Points", "bullets": ["Un", "Deux"]},
            ],
        }
    )
    assert deck["document"]["filename"].endswith(".pptx")

    # Missing essentials are refused rather than producing an empty file.
    assert tools.get("office_document").handler({"title": "X"})["error"] == "title_required"
    assert (
        tools.get("office_document").handler({"kind": "report", "title": "Rapport"})["error"]
        == "blocks_required"
    )
    assert tools.get("office_spreadsheet").handler({"title": "Vide"})["error"] == "sheets_required"
    assert tools.get("office_deck").handler({"title": "Vide"})["error"] == "slides_required"


def test_a_produced_workbook_can_be_shown_as_a_card(tmp_path):
    """The Livrables list and the visual cards work for every format, not just
    Word."""
    app = build_app(tmp_path)
    tools = app.state.agent.tools
    workbook = tools.get("office_spreadsheet").handler(
        {
            "title": "Budget",
            "sheets": [{"name": "B", "columns": [{"header": "Poste"}], "rows": [["Loyer"]]}],
        }
    )
    from emefa.domain.visuals import CardCollector

    with CardCollector() as collector:
        shown = tools.get("show_document").handler(
            {"document_id": workbook["document"]["document_id"]}
        )
    assert shown["shown"] is True
    assert collector.cards[0].title == "Budget"
