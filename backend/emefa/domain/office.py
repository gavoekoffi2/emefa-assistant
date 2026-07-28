"""Office capability boundary: structured specs, then a provider adapter.

CLAUDE.md §19 requires business logic to depend on a *capability* interface
rather than on a concrete office engine. The layering is therefore:

    caller (skills / workflows)
        -> DocumentSpec | WorkbookSpec | DeckSpec  (this module, engine-free)
            -> OfficeProvider protocol
                -> PythonOfficeProvider (python-docx / openpyxl / python-pptx)

Replacing the engine (OfficeCLI, a service, a template renderer) means writing
one new provider, never touching a caller.

Two product rules are encoded here rather than left to the engine:
  * generated files stay *editable* — no flattening to PDF, no images of text;
  * spreadsheet formulas stay *live* — a cell whose value starts with ``=`` is
    written as a real formula, and totals rows are generated as ``SUM`` ranges,
    so the recipient can change an input and watch the sheet recompute.
"""

from __future__ import annotations

import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Literal, Protocol

DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
XLSX_MIME = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

ArtifactKind = Literal["document", "workbook", "presentation"]

KIND_EXTENSIONS: dict[str, str] = {
    "document": ".docx",
    "workbook": ".xlsx",
    "presentation": ".pptx",
}
KIND_MIME_TYPES: dict[str, str] = {
    "document": DOCX_MIME,
    "workbook": XLSX_MIME,
    "presentation": PPTX_MIME,
}

# EMEFA brand accents, kept subtle enough for a business document.
_ACCENT = "0B6E8F"
_INK = "1B2733"
_MUTED = "5A6B7B"


# --------------------------------------------------------------------------
# Specs — plain data, no engine import
# --------------------------------------------------------------------------


@dataclass(frozen=True, slots=True)
class Block:
    """One structural element of a written document."""

    kind: Literal["heading", "paragraph", "bullet", "numbered", "table", "spacer"]
    text: str = ""
    level: int = 1
    rows: tuple[tuple[str, ...], ...] = ()


@dataclass(frozen=True, slots=True)
class DocumentSpec:
    title: str
    subtitle: str = ""
    blocks: tuple[Block, ...] = ()
    footer: str = ""


@dataclass(frozen=True, slots=True)
class SheetSpec:
    name: str
    columns: tuple[str, ...] = ()
    rows: tuple[tuple[Any, ...], ...] = ()
    #: Column letters (or 1-based indexes) that should receive a SUM total row.
    total_columns: tuple[str, ...] = ()
    number_format: str = "#,##0.00"
    notes: str = ""


@dataclass(frozen=True, slots=True)
class WorkbookSpec:
    title: str
    sheets: tuple[SheetSpec, ...] = ()


@dataclass(frozen=True, slots=True)
class SlideSpec:
    title: str
    bullets: tuple[str, ...] = ()
    notes: str = ""


@dataclass(frozen=True, slots=True)
class DeckSpec:
    title: str
    subtitle: str = ""
    slides: tuple[SlideSpec, ...] = field(default_factory=tuple)


class OfficeProvider(Protocol):
    """Everything EMEFA needs from an office engine."""

    def render_document(self, spec: DocumentSpec, destination: Path) -> None: ...

    def render_workbook(self, spec: WorkbookSpec, destination: Path) -> None: ...

    def render_presentation(self, spec: DeckSpec, destination: Path) -> None: ...

    def read_document(self, source: Path) -> tuple[str, str]:
        """Return ``(title, body_text)`` for an existing document."""
        ...


# --------------------------------------------------------------------------
# Content parsing — a light, forgiving structure syntax
# --------------------------------------------------------------------------

_HEADING = re.compile(r"^(#{1,3})\s+(.*)$")
_BULLET = re.compile(r"^[-*•]\s+(.*)$")
_NUMBERED = re.compile(r"^\d+[.)]\s+(.*)$")
_TABLE_SEPARATOR = re.compile(r"^\|[\s:|-]+\|$")


def _table_cells(line: str) -> tuple[str, ...]:
    return tuple(cell.strip() for cell in line.strip().strip("|").split("|"))


def parse_content(content: str) -> tuple[Block, ...]:
    """Turn assistant-written text into document structure.

    The syntax is deliberately the subset an LLM produces naturally:
    ``#``/``##``/``###`` headings, ``-`` bullets, ``1.`` numbered items and
    ``|`` tables. Anything else is a paragraph, so plain prose still works
    exactly as before.
    """
    blocks: list[Block] = []
    lines = str(content).replace("\r\n", "\n").split("\n")
    index = 0
    while index < len(lines):
        raw = lines[index]
        line = raw.strip()
        index += 1
        if not line:
            if blocks and blocks[-1].kind != "spacer":
                blocks.append(Block(kind="spacer"))
            continue
        heading = _HEADING.match(line)
        if heading is not None:
            blocks.append(
                Block(kind="heading", text=heading.group(2).strip(), level=len(heading.group(1)))
            )
            continue
        if line.startswith("|") and line.endswith("|"):
            rows: list[tuple[str, ...]] = [_table_cells(line)]
            while index < len(lines):
                candidate = lines[index].strip()
                if not (candidate.startswith("|") and candidate.endswith("|")):
                    break
                index += 1
                if _TABLE_SEPARATOR.match(candidate):
                    continue
                rows.append(_table_cells(candidate))
            blocks.append(Block(kind="table", rows=tuple(rows)))
            continue
        bullet = _BULLET.match(line)
        if bullet is not None:
            blocks.append(Block(kind="bullet", text=bullet.group(1).strip()))
            continue
        numbered = _NUMBERED.match(line)
        if numbered is not None:
            blocks.append(Block(kind="numbered", text=numbered.group(1).strip()))
            continue
        blocks.append(Block(kind="paragraph", text=line))
    while blocks and blocks[-1].kind == "spacer":
        blocks.pop()
    return tuple(blocks)


def blocks_to_text(blocks: tuple[Block, ...]) -> str:
    """Inverse of :func:`parse_content`, used when re-reading an artifact."""
    lines: list[str] = []
    for block in blocks:
        if block.kind == "spacer":
            lines.append("")
        elif block.kind == "heading":
            lines.append(f"{'#' * max(1, min(block.level, 3))} {block.text}")
        elif block.kind == "bullet":
            lines.append(f"- {block.text}")
        elif block.kind == "numbered":
            lines.append(f"1. {block.text}")
        elif block.kind == "table":
            lines.extend("| " + " | ".join(row) + " |" for row in block.rows)
        else:
            lines.append(block.text)
    return "\n".join(lines)


# --------------------------------------------------------------------------
# Default provider
# --------------------------------------------------------------------------


class PythonOfficeProvider:
    """Adapter over python-docx, openpyxl and python-pptx.

    Imports are deferred to call time so the FastAPI app still boots (and the
    document capability degrades to a clear error) if an optional engine is
    missing from the deployment image.
    """

    # -- Word ------------------------------------------------------------
    def render_document(self, spec: DocumentSpec, destination: Path) -> None:
        from docx import Document
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.shared import Pt, RGBColor

        document = Document()
        base = document.styles["Normal"]
        base.font.name = "Calibri"
        base.font.size = Pt(11)
        base.paragraph_format.space_after = Pt(6)
        base.paragraph_format.line_spacing = 1.15

        heading = document.add_heading(spec.title, level=0)
        for run in heading.runs:
            run.font.color.rgb = RGBColor.from_string(_INK)
        if spec.subtitle:
            subtitle = document.add_paragraph(spec.subtitle)
            subtitle.alignment = WD_ALIGN_PARAGRAPH.LEFT
            for run in subtitle.runs:
                run.italic = True
                run.font.size = Pt(10)
                run.font.color.rgb = RGBColor.from_string(_MUTED)

        for block in spec.blocks:
            if block.kind == "spacer":
                continue
            if block.kind == "heading":
                level = max(1, min(block.level, 3))
                written = document.add_heading(block.text, level=level)
                for run in written.runs:
                    run.font.color.rgb = RGBColor.from_string(_ACCENT if level == 1 else _INK)
            elif block.kind == "bullet":
                document.add_paragraph(block.text, style="List Bullet")
            elif block.kind == "numbered":
                document.add_paragraph(block.text, style="List Number")
            elif block.kind == "table" and block.rows:
                width = max(len(row) for row in block.rows)
                table = document.add_table(rows=0, cols=width)
                table.style = "Light Grid Accent 1"
                for position, row in enumerate(block.rows):
                    cells = table.add_row().cells
                    for column in range(width):
                        cells[column].text = row[column] if column < len(row) else ""
                        if position == 0:
                            for paragraph in cells[column].paragraphs:
                                for run in paragraph.runs:
                                    run.bold = True
            else:
                document.add_paragraph(block.text)

        if spec.footer:
            section = document.sections[0]
            paragraph = section.footer.paragraphs[0]
            paragraph.text = spec.footer
            paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
            for run in paragraph.runs:
                run.font.size = Pt(8)
                run.font.color.rgb = RGBColor.from_string(_MUTED)

        self._atomic_save(document.save, destination)

    def read_document(self, source: Path) -> tuple[str, str]:
        from docx import Document

        document = Document(str(source))
        paragraphs = list(document.paragraphs)
        title = paragraphs[0].text if paragraphs else "Document EMEFA"
        body: list[str] = []
        for paragraph in paragraphs[1:]:
            style = (paragraph.style.name or "").lower()
            text = paragraph.text
            if not text.strip():
                body.append("")
            elif style.startswith("heading"):
                digits = "".join(char for char in style if char.isdigit())
                body.append(f"{'#' * max(1, min(int(digits or 1), 3))} {text}")
            elif style.startswith("list bullet"):
                body.append(f"- {text}")
            elif style.startswith("list number"):
                body.append(f"1. {text}")
            else:
                body.append(text)
        for table in document.tables:
            for row in table.rows:
                body.append("| " + " | ".join(cell.text for cell in row.cells) + " |")
        return title, "\n".join(body).strip()

    # -- Excel -----------------------------------------------------------
    def render_workbook(self, spec: WorkbookSpec, destination: Path) -> None:
        from openpyxl import Workbook
        from openpyxl.styles import Alignment, Font, PatternFill
        from openpyxl.utils import get_column_letter

        workbook = Workbook()
        workbook.remove(workbook.active)
        header_fill = PatternFill("solid", fgColor=_ACCENT)
        header_font = Font(bold=True, color="FFFFFF")

        for sheet_spec in spec.sheets or (SheetSpec(name="Feuille 1"),):
            sheet = workbook.create_sheet(title=sheet_spec.name[:31] or "Feuille")
            row_index = 1
            if sheet_spec.columns:
                for column_index, label in enumerate(sheet_spec.columns, start=1):
                    cell = sheet.cell(row=1, column=column_index, value=label)
                    cell.fill = header_fill
                    cell.font = header_font
                    cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
                sheet.freeze_panes = "A2"
                row_index = 2

            first_data_row = row_index
            for row in sheet_spec.rows:
                for column_index, value in enumerate(row, start=1):
                    cell = sheet.cell(row=row_index, column=column_index, value=value)
                    # openpyxl stores a leading "=" string as a live formula.
                    if isinstance(value, (int, float)) and not isinstance(value, bool):
                        cell.number_format = sheet_spec.number_format
                row_index += 1
            last_data_row = row_index - 1

            if sheet_spec.total_columns and last_data_row >= first_data_row:
                label_cell = sheet.cell(row=row_index, column=1, value="Total")
                label_cell.font = Font(bold=True)
                for reference in sheet_spec.total_columns:
                    letter = self._column_letter(reference)
                    if letter is None:
                        continue
                    total = sheet[f"{letter}{row_index}"]
                    total.value = f"=SUM({letter}{first_data_row}:{letter}{last_data_row})"
                    total.font = Font(bold=True)
                    total.number_format = sheet_spec.number_format
                row_index += 1

            if sheet_spec.notes:
                sheet.cell(row=row_index + 1, column=1, value=sheet_spec.notes)

            column_count = max(
                len(sheet_spec.columns),
                max((len(row) for row in sheet_spec.rows), default=0),
                1,
            )
            for column_index in range(1, column_count + 1):
                letter = get_column_letter(column_index)
                longest = max(
                    (
                        len(str(sheet[f"{letter}{index}"].value or ""))
                        for index in range(1, row_index + 1)
                    ),
                    default=10,
                )
                sheet.column_dimensions[letter].width = min(46, max(12, longest + 3))
            if sheet_spec.columns and last_data_row >= first_data_row:
                sheet.auto_filter.ref = (
                    f"A1:{get_column_letter(len(sheet_spec.columns))}{last_data_row}"
                )

        self._atomic_save(workbook.save, destination)

    @staticmethod
    def _column_letter(reference: str) -> str | None:
        from openpyxl.utils import get_column_letter

        value = str(reference).strip()
        if value.isdigit():
            index = int(value)
            return get_column_letter(index) if 1 <= index <= 16_384 else None
        return value.upper() if value.isalpha() and len(value) <= 3 else None

    # -- PowerPoint ------------------------------------------------------
    def render_presentation(self, spec: DeckSpec, destination: Path) -> None:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.util import Pt

        presentation = Presentation()
        cover = presentation.slides.add_slide(presentation.slide_layouts[0])
        cover.shapes.title.text = spec.title
        if len(cover.placeholders) > 1:
            cover.placeholders[1].text = spec.subtitle

        for slide_spec in spec.slides:
            layout = presentation.slide_layouts[1 if slide_spec.bullets else 5]
            slide = presentation.slides.add_slide(layout)
            slide.shapes.title.text = slide_spec.title
            for paragraph in slide.shapes.title.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor.from_string(_INK)
                    run.font.size = Pt(30)
            if slide_spec.bullets:
                body = slide.placeholders[1].text_frame
                body.clear()
                for position, bullet in enumerate(slide_spec.bullets):
                    paragraph = body.paragraphs[0] if position == 0 else body.add_paragraph()
                    paragraph.text = bullet
                    paragraph.level = 0
                    for run in paragraph.runs:
                        run.font.size = Pt(18)
            if slide_spec.notes:
                slide.notes_slide.notes_text_frame.text = slide_spec.notes

        self._atomic_save(presentation.save, destination)

    # -- shared ----------------------------------------------------------
    @staticmethod
    def _atomic_save(save: Any, destination: Path) -> None:
        import os

        temporary = destination.with_name(destination.name + ".tmp")
        save(str(temporary))
        os.replace(temporary, destination)
