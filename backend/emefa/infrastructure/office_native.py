"""Native office renderer — python-docx, openpyxl, python-pptx.

The one design rule running through this file: **a spreadsheet must stay a
spreadsheet.** Writing computed values where a formula belongs produces a file
that looks right and dies the moment someone changes a quantity. So totals and
line amounts are written as live formulas, and the same figures are *also*
computed in Python and returned in `BuiltArtifact.computed` — because the
assistant has to be able to say "le total est de 4 720 000 FCFA" without
opening the file, and must not invent that number either.

The Word table of contents is a field, not a rendered list. Word fills it in
when the document is opened and fields are updated, which is how every real
template does it; a hand-rendered list of headings goes stale the first time
anyone edits the document.
"""

from __future__ import annotations

import io
import re
from datetime import date
from typing import Any

from docx import Document
from docx.enum.section import WD_SECTION
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Pt, RGBColor
from openpyxl import Workbook
from openpyxl.chart import BarChart, LineChart, PieChart, Reference
from openpyxl.styles import Alignment, Font, PatternFill
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor as PptxColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt as PptxPt

from emefa.domain.office.provider import OfficeError
from emefa.domain.office.schemas import (
    MAX_BLOCKS,
    MAX_TABLE_ROWS,
    Block,
    BlockKind,
    BuiltArtifact,
    Chart,
    ChartKind,
    ColumnFormat,
    DeckSpec,
    DocumentSpec,
    Sheet,
    Slide,
    SlideLayout,
    WorkbookSpec,
)

_NUMBER_FORMATS = {
    ColumnFormat.TEXT: "@",
    ColumnFormat.NUMBER: "#,##0.##",
    ColumnFormat.MONEY: "#,##0",
    ColumnFormat.PERCENT: "0.0%",
    ColumnFormat.DATE: "yyyy-mm-dd",
}

_SAFE_FORMULA_FUNCTIONS = {"ABS", "IF", "MAX", "MIN", "ROUND", "SUM"}
_FORMULA_TOKEN = re.compile(
    r"\s+|\$?[A-Z]{1,3}\$?(?:\d+|\{row\})|\d+(?:\.\d+)?|"
    r"[A-Z_][A-Z0-9_]*|<=|>=|<>|[=+\-*/^%(),:<>]",
    re.IGNORECASE,
)
_CELL_TOKEN = re.compile(r"\$?[A-Z]{1,3}\$?(?:\d+|\{row\})", re.IGNORECASE)


def validate_local_formula(formula: str) -> str:
    """Allow only local cell arithmetic and a small safe function subset."""
    value = formula.strip()
    if not value.startswith("=") or len(value) > 200:
        raise OfficeError("formule Excel refusée : syntaxe non sûre")
    position = 0
    words: list[tuple[str, int]] = []
    for match in _FORMULA_TOKEN.finditer(value):
        if match.start() != position:
            raise OfficeError("formule Excel refusée : référence externe ou caractère interdit")
        token = match.group(0)
        if (
            re.fullmatch(r"[A-Z_][A-Z0-9_]*", token, re.IGNORECASE)
            and not _CELL_TOKEN.fullmatch(token)
        ):
            words.append((token.upper(), match.end()))
        position = match.end()
    if position != len(value):
        raise OfficeError("formule Excel refusée : référence externe ou caractère interdit")
    for word, end in words:
        following = value[end:].lstrip()
        if word not in _SAFE_FORMULA_FUNCTIONS or not following.startswith("("):
            raise OfficeError("formule Excel refusée : fonction non autorisée")
    return value


def _rgb(value: str) -> RGBColor:
    cleaned = re.sub(r"[^0-9A-Fa-f]", "", value or "")[:6].ljust(6, "0")
    return RGBColor.from_string(cleaned.upper())


def _hex(value: str) -> str:
    return re.sub(r"[^0-9A-Fa-f]", "", value or "")[:6].ljust(6, "0").upper()


def _field(paragraph: Any, instruction: str) -> None:
    """Insert a Word field (PAGE, TOC).

    A field is a live instruction Word evaluates; the alternative — writing the
    page number as text — is wrong on every page but one.
    """
    run = paragraph.add_run()
    begin = OxmlElement("w:fldChar")
    begin.set(qn("w:fldCharType"), "begin")
    instruction_element = OxmlElement("w:instrText")
    instruction_element.set(qn("xml:space"), "preserve")
    instruction_element.text = instruction
    separate = OxmlElement("w:fldChar")
    separate.set(qn("w:fldCharType"), "separate")
    end = OxmlElement("w:fldChar")
    end.set(qn("w:fldCharType"), "end")
    for element in (begin, instruction_element, separate, end):
        run._r.append(element)


class NativeOfficeProvider:
    """Default renderer. No external process, no network, no licence."""

    name = "native"

    # ── Word ──────────────────────────────────────────────────────────────

    def build_document(self, spec: DocumentSpec) -> BuiltArtifact:
        if not spec.title.strip():
            raise OfficeError("un document sans titre n'est pas un document")

        document = Document()
        warnings: list[str] = []
        brand = spec.brand

        self._apply_base_style(document, brand)
        self._apply_header_footer(document, spec)
        self._write_cover(document, spec)

        if spec.table_of_contents:
            document.add_heading("Sommaire", level=1)
            _field(
                document.add_paragraph(),
                r'TOC \o "1-3" \h \z \u',
            )
            # Said out loud, because a user who opens the file and sees an
            # empty heading will otherwise think it is broken.
            warnings.append(
                "Le sommaire se remplit à l'ouverture du document, quand Word "
                "met les champs à jour (Ctrl+A puis F9 pour forcer)."
            )
            document.add_page_break()

        blocks = spec.blocks[:MAX_BLOCKS]
        if len(spec.blocks) > MAX_BLOCKS:
            warnings.append(f"Document tronqué à {MAX_BLOCKS} blocs.")
        for block in blocks:
            self._write_block(document, block, brand)

        buffer = io.BytesIO()
        document.save(buffer)
        return BuiltArtifact(buffer.getvalue(), "docx", {}, tuple(warnings))

    def _apply_base_style(self, document: Any, brand: Any) -> None:
        normal = document.styles["Normal"]
        normal.font.name = brand.font or "Calibri"
        normal.font.size = Pt(11)
        for level in range(1, 4):
            style = document.styles[f"Heading {level}"]
            style.font.color.rgb = _rgb(brand.primary_colour)
            style.font.name = brand.font or "Calibri"

    def _apply_header_footer(self, document: Any, spec: DocumentSpec) -> None:
        section = document.sections[0]
        header_text = spec.header_text or spec.brand.company_name
        if header_text:
            paragraph = section.header.paragraphs[0]
            paragraph.text = header_text
            paragraph.alignment = WD_ALIGN_PARAGRAPH.RIGHT
            paragraph.runs[0].font.size = Pt(8)

        footer_bits = [spec.footer_text or spec.brand.footer_note, spec.brand.contact]
        footer = section.footer.paragraphs[0]
        footer.text = " · ".join(bit for bit in footer_bits if bit)
        footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
        if footer.runs:
            footer.runs[0].font.size = Pt(8)
        if spec.page_numbers:
            if footer.text:
                footer.add_run("  —  ").font.size = Pt(8)
            _field(footer, "PAGE")

    def _write_cover(self, document: Any, spec: DocumentSpec) -> None:
        brand = spec.brand
        if brand.company_name:
            heading = document.add_paragraph(brand.company_name)
            heading.runs[0].bold = True
            heading.runs[0].font.color.rgb = _rgb(brand.primary_colour)
        if brand.address:
            small = document.add_paragraph(brand.address)
            small.runs[0].font.size = Pt(9)

        document.add_heading(spec.title, level=0)
        if spec.subtitle:
            subtitle = document.add_paragraph(spec.subtitle)
            subtitle.runs[0].italic = True

        details = [("Référence", spec.reference), ("Date", spec.document_date)]
        if spec.recipient:
            details.insert(0, ("Destinataire", spec.recipient))
        line = " · ".join(f"{label} : {value}" for label, value in details if value)
        if line:
            marker = document.add_paragraph(line)
            marker.runs[0].font.size = Pt(9)

    def _write_block(self, document: Any, block: Block, brand: Any) -> None:
        if block.kind is BlockKind.HEADING:
            document.add_heading(block.text, level=max(1, min(4, block.level)))
        elif block.kind is BlockKind.PARAGRAPH:
            document.add_paragraph(block.text[:20_000])
        elif block.kind is BlockKind.BULLETS:
            for item in block.items:
                document.add_paragraph(item, style="List Bullet")
        elif block.kind is BlockKind.NUMBERED:
            for item in block.items:
                document.add_paragraph(item, style="List Number")
        elif block.kind is BlockKind.FIELDS:
            table = document.add_table(rows=0, cols=2)
            table.style = "Table Grid"
            for label, value in block.fields:
                cells = table.add_row().cells
                cells[0].text = label
                cells[0].paragraphs[0].runs[0].bold = True
                cells[1].text = value
        elif block.kind is BlockKind.TABLE:
            self._write_table(document, block, brand)
        elif block.kind is BlockKind.SIGNATURE:
            document.add_paragraph()
            table = document.add_table(rows=1, cols=2)
            for index, label in enumerate(block.items[:2] or ("Le client", "L'entreprise")):
                cell = table.cell(0, index)
                cell.text = f"{label}\n\nNom :\nDate :\nSignature :"
        elif block.kind is BlockKind.PAGE_BREAK:
            document.add_page_break()

    def _write_table(self, document: Any, block: Block, brand: Any) -> None:
        if not block.columns:
            return
        table = document.add_table(rows=1, cols=len(block.columns))
        table.style = "Table Grid"
        for index, header in enumerate(block.columns):
            cell = table.rows[0].cells[index]
            cell.text = header
            run = cell.paragraphs[0].runs[0]
            run.bold = True
            run.font.color.rgb = _rgb(brand.primary_colour)
        for row in block.rows[:MAX_TABLE_ROWS]:
            cells = table.add_row().cells
            for index, value in enumerate(row[: len(block.columns)]):
                cells[index].text = str(value)
                if block.numeric_last_column and index == len(block.columns) - 1:
                    cells[index].paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.RIGHT

    # ── Excel ─────────────────────────────────────────────────────────────

    def build_workbook(self, spec: WorkbookSpec) -> BuiltArtifact:
        if not spec.sheets:
            raise OfficeError("un classeur sans feuille n'a rien à montrer")
        for sheet in spec.sheets:
            for column in sheet.columns:
                if column.formula:
                    validate_local_formula(column.formula)

        workbook = Workbook()
        workbook.remove(workbook.active)
        computed: dict[str, float] = {}
        warnings: list[str] = []

        for sheet_spec in spec.sheets:
            self._write_sheet(workbook, sheet_spec, spec.brand, computed, warnings)

        buffer = io.BytesIO()
        workbook.save(buffer)
        return BuiltArtifact(buffer.getvalue(), "xlsx", computed, tuple(warnings))

    def _write_sheet(
        self,
        workbook: Any,
        sheet: Sheet,
        brand: Any,
        computed: dict[str, float],
        warnings: list[str],
    ) -> None:
        # Excel refuses these characters in a tab name, and refuses silently
        # enough that the failure surfaces much later.
        safe_name = re.sub(r"[\\/*?:\[\]]", "-", sheet.name or "Feuille")[:31]
        worksheet = workbook.create_sheet(safe_name)
        if not sheet.columns:
            return

        first_row = 1
        if sheet.caption:
            worksheet.cell(row=1, column=1, value=sheet.caption).font = Font(italic=True, size=9)
            first_row = 3

        header_fill = PatternFill("solid", fgColor=_hex(brand.primary_colour))
        for index, column in enumerate(sheet.columns, start=1):
            cell = worksheet.cell(row=first_row, column=index, value=column.header)
            cell.font = Font(bold=True, color="FFFFFF")
            cell.fill = header_fill
            cell.alignment = Alignment(horizontal="center", vertical="center")
            worksheet.column_dimensions[get_column_letter(index)].width = (
                column.width or max(12, min(40, len(column.header) + 6))
            )

        rows = sheet.rows[:MAX_TABLE_ROWS]
        if len(sheet.rows) > MAX_TABLE_ROWS:
            warnings.append(f"Feuille « {safe_name} » tronquée à {MAX_TABLE_ROWS} lignes.")

        for offset, row in enumerate(rows):
            excel_row = first_row + 1 + offset
            for index, column in enumerate(sheet.columns, start=1):
                cell = worksheet.cell(row=excel_row, column=index)
                if column.formula:
                    # Live formula, so the sheet recalculates when someone
                    # changes a quantity. This is the difference between a
                    # spreadsheet and a picture of one.
                    cell.value = column.formula.replace("{row}", str(excel_row))
                else:
                    cell.value = row[index - 1] if index - 1 < len(row) else None
                cell.number_format = _NUMBER_FORMATS[column.format]

        total_row = first_row + len(rows) + 1
        for index, column in enumerate(sheet.columns, start=1):
            if not column.total:
                continue
            letter = get_column_letter(index)
            label_cell = worksheet.cell(row=total_row, column=max(1, index - 1))
            if index > 1 and not label_cell.value:
                label_cell.value = "Total"
                label_cell.font = Font(bold=True)
            cell = worksheet.cell(row=total_row, column=index)
            cell.value = f"=SUM({letter}{first_row + 1}:{letter}{first_row + len(rows)})"
            cell.font = Font(bold=True)
            cell.number_format = _NUMBER_FORMATS[column.format]
            # The same sum, computed here, so the assistant can state the
            # figure without opening the file and without guessing it.
            #
            # A formula column holds no literal values, so summing the cells
            # would report zero — and an assistant announcing "total : 0" on a
            # quote is worse than one that says it cannot tell yet.
            total = _column_total(rows, sheet.columns, index - 1)
            if total is None:
                warnings.append(
                    f"Le total de « {column.header} » apparaîtra à l'ouverture du "
                    "fichier : la formule est trop complexe pour être calculée ici."
                )
            else:
                computed[f"{safe_name}.{column.header}"] = float(total)

        if sheet.freeze_header:
            worksheet.freeze_panes = worksheet.cell(row=first_row + 1, column=1)
        if sheet.autofilter and rows:
            worksheet.auto_filter.ref = (
                f"A{first_row}:{get_column_letter(len(sheet.columns))}"
                f"{first_row + len(rows)}"
            )
        self._add_sheet_chart(worksheet, sheet, first_row, len(rows), warnings)

    def _add_sheet_chart(
        self, worksheet: Any, sheet: Sheet, header_row: int, row_count: int, warnings: list[str]
    ) -> None:
        chart_spec = sheet.chart
        if chart_spec.kind is ChartKind.NONE or row_count == 0:
            return
        headers = [column.header for column in sheet.columns]
        if chart_spec.label_column not in headers:
            warnings.append("Graphique ignoré : colonne d'étiquettes introuvable.")
            return
        value_indexes = [
            headers.index(name) + 1 for name in chart_spec.value_columns if name in headers
        ]
        if not value_indexes:
            warnings.append("Graphique ignoré : aucune colonne de valeurs reconnue.")
            return

        chart = {
            ChartKind.BAR: BarChart,
            ChartKind.LINE: LineChart,
            ChartKind.PIE: PieChart,
        }[chart_spec.kind]()
        chart.title = chart_spec.title or sheet.name
        labels = Reference(
            worksheet,
            min_col=headers.index(chart_spec.label_column) + 1,
            min_row=header_row + 1,
            max_row=header_row + row_count,
        )
        for column_index in value_indexes:
            data = Reference(
                worksheet,
                min_col=column_index,
                min_row=header_row,
                max_row=header_row + row_count,
            )
            chart.add_data(data, titles_from_data=True)
        chart.set_categories(labels)
        chart.height, chart.width = 8, 16
        worksheet.add_chart(chart, f"{get_column_letter(len(sheet.columns) + 2)}{header_row}")

    # ── PowerPoint ────────────────────────────────────────────────────────

    def build_deck(self, spec: DeckSpec) -> BuiltArtifact:
        if not spec.slides:
            raise OfficeError("une présentation sans diapositive n'existe pas")

        presentation = Presentation()
        warnings: list[str] = []
        primary = PptxColor.from_string(_hex(spec.brand.primary_colour))

        for slide_spec in spec.slides[:60]:
            self._write_slide(presentation, slide_spec, spec, primary, warnings)

        buffer = io.BytesIO()
        presentation.save(buffer)
        return BuiltArtifact(buffer.getvalue(), "pptx", {}, tuple(warnings))

    def _write_slide(
        self, presentation: Any, slide_spec: Slide, spec: DeckSpec, primary: Any, warnings: list[str]
    ) -> None:
        layouts = presentation.slide_layouts
        if slide_spec.layout in {SlideLayout.TITLE, SlideLayout.CLOSING}:
            slide = presentation.slides.add_slide(layouts[0])
            slide.shapes.title.text = slide_spec.title or spec.title
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = slide_spec.subtitle or spec.subtitle
        elif slide_spec.layout is SlideLayout.SECTION:
            slide = presentation.slides.add_slide(layouts[2])
            slide.shapes.title.text = slide_spec.title
        elif slide_spec.layout is SlideLayout.QUOTE:
            slide = presentation.slides.add_slide(layouts[5])
            slide.shapes.title.text = slide_spec.title
            box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
            paragraph = box.text_frame.paragraphs[0]
            paragraph.text = slide_spec.subtitle
            paragraph.font.size = PptxPt(24)
            paragraph.font.italic = True
        elif slide_spec.layout is SlideLayout.TABLE:
            slide = presentation.slides.add_slide(layouts[5])
            slide.shapes.title.text = slide_spec.title
            self._add_slide_table(slide, slide_spec, primary, warnings)
        elif slide_spec.layout is SlideLayout.CHART:
            slide = presentation.slides.add_slide(layouts[5])
            slide.shapes.title.text = slide_spec.title
            self._add_slide_chart(slide, slide_spec, warnings)
        else:
            slide = presentation.slides.add_slide(layouts[1])
            slide.shapes.title.text = slide_spec.title
            frame = slide.placeholders[1].text_frame
            frame.clear()
            for index, bullet in enumerate(slide_spec.bullets[:8]):
                paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
                paragraph.text = bullet
                paragraph.level = 0

        if slide.shapes.title is not None and slide.shapes.title.text_frame.paragraphs:
            slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = primary
        if slide_spec.notes:
            slide.notes_slide.notes_text_frame.text = slide_spec.notes

    def _add_slide_table(self, slide: Any, slide_spec: Slide, primary: Any, warnings: list[str]) -> None:
        if not slide_spec.columns:
            warnings.append(f"Diapositive « {slide_spec.title} » : tableau sans colonnes.")
            return
        rows = slide_spec.rows[:10]
        if len(slide_spec.rows) > 10:
            warnings.append(
                f"Diapositive « {slide_spec.title} » : tableau réduit à 10 lignes lisibles."
            )
        shape = slide.shapes.add_table(
            len(rows) + 1, len(slide_spec.columns),
            Inches(0.6), Inches(1.8), Inches(9), Inches(0.4 * (len(rows) + 1)),
        )
        table = shape.table
        for index, header in enumerate(slide_spec.columns):
            table.cell(0, index).text = header
        for row_index, row in enumerate(rows, start=1):
            for column_index in range(len(slide_spec.columns)):
                value = row[column_index] if column_index < len(row) else ""
                table.cell(row_index, column_index).text = str(value)

    def _add_slide_chart(self, slide: Any, slide_spec: Slide, warnings: list[str]) -> None:
        points = [(label, value) for label, value in slide_spec.points if str(label).strip()]
        if not points:
            warnings.append(f"Diapositive « {slide_spec.title} » : graphique sans données.")
            return
        data = CategoryChartData()
        data.categories = [label for label, _ in points]
        data.add_series(slide_spec.chart.title or "Série", [float(value) for _, value in points])
        chart_type = {
            ChartKind.BAR: XL_CHART_TYPE.COLUMN_CLUSTERED,
            ChartKind.LINE: XL_CHART_TYPE.LINE_MARKERS,
            ChartKind.PIE: XL_CHART_TYPE.PIE,
            ChartKind.NONE: XL_CHART_TYPE.COLUMN_CLUSTERED,
        }[slide_spec.chart.kind]
        slide.shapes.add_chart(
            chart_type, Inches(1), Inches(1.8), Inches(8), Inches(4.5), data
        )


def _numeric(row: tuple[Any, ...], index: int) -> float | None:
    if index >= len(row):
        return None
    value = row[index]
    if isinstance(value, bool):
        return None
    return float(value) if isinstance(value, (int, float)) else None


#: A row formula made of cell references, numbers and the four operators.
#: Anything richer is left to Excel rather than half-evaluated here.
_TOKEN = re.compile(r"([A-Z]{1,3})\{row\}|(\d+(?:\.\d+)?)|([+\-*/])")


def _evaluate_row_formula(formula: str, row: tuple[Any, ...]) -> float | None:
    """Evaluate a simple per-row formula against literal cell values.

    Not an expression engine, and deliberately not `eval`. It handles the
    shapes an assistant actually writes — `=B{row}*C{row}`,
    `=D{row}*1.18`, `=B{row}-C{row}` — and returns None for anything else so
    the caller says "the total appears when you open it" instead of reporting
    a number it half-guessed.
    """
    body = formula.strip().lstrip("=")
    if not body or any(character in body for character in "():,$"):
        return None

    values: list[float] = []
    operators: list[str] = []
    position = 0
    for match in _TOKEN.finditer(body):
        if match.start() != position:
            return None  # something unparsed between tokens
        position = match.end()
        reference, number, operator = match.groups()
        if reference is not None:
            index = 0
            for character in reference:
                index = index * 26 + (ord(character) - 64)
            value = _numeric(row, index - 1)
            if value is None:
                return None
            values.append(value)
        elif number is not None:
            values.append(float(number))
        else:
            operators.append(operator)
    if position != len(body) or len(values) != len(operators) + 1:
        return None

    # Multiplication and division first, then the rest left to right.
    for pass_operators in (("*", "/"), ("+", "-")):
        index = 0
        while index < len(operators):
            if operators[index] not in pass_operators:
                index += 1
                continue
            left, right = values[index], values[index + 1]
            operator = operators.pop(index)
            if operator == "*":
                result = left * right
            elif operator == "/":
                if right == 0:
                    return None
                result = left / right
            elif operator == "+":
                result = left + right
            else:
                result = left - right
            values[index : index + 2] = [result]
    return values[0] if len(values) == 1 else None


def _column_total(
    rows: tuple[tuple[Any, ...], ...], columns: tuple[Any, ...], index: int
) -> float | None:
    """Sum a column the way Excel will, or admit it cannot."""
    formula = columns[index].formula if index < len(columns) else ""
    total = 0.0
    for row in rows:
        value = (
            _evaluate_row_formula(formula, row) if formula else _numeric(row, index)
        )
        if value is None:
            # A blank cell in a literal column is a zero; an unevaluable
            # formula is not.
            if formula:
                return None
            continue
        total += value
    return total
